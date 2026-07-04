#!/usr/bin/env python3
"""office_audit.py — checagens de acessibilidade em docx/pptx/xlsx.

O checker do MS Office não tem CLI nem API: as regras são públicas, o motor é
preso à GUI. Este script reimplementa o subconjunto verificável por máquina
sobre o OOXML. O que exige julgamento (qualidade do alt-text, ordem de leitura
lógica) sai como item para triagem LLM/humana, nunca como veredito automático.

Requer extras: pip install python-docx python-pptx openpyxl

Uso: python3 office_audit.py <arquivo.docx|pptx|xlsx> [--json]

Taxonomia (espelha o checker da Microsoft): erro | aviso | dica.
Cada achado sai com proveniência "motor" (este script) e critério WCAG.
"""

import json
import sys
from pathlib import Path


def achado(sev, regra, wcag, local, detalhe):
    return {"severidade": sev, "regra": regra, "wcag": wcag,
            "local": local, "detalhe": detalhe, "proveniencia": "motor"}


def _descr(elemento):
    """Alt-text no OOXML vive no atributo descr de cNvPr/docPr (sem API nas libs)."""
    nos = elemento.xpath('.//*[local-name()="cNvPr" or local-name()="docPr"]')
    return (nos[0].get("descr") or "").strip() if nos else ""


def auditar_docx(caminho):
    import docx
    doc = docx.Document(caminho)
    achados = []
    if not (doc.core_properties.title or "").strip():
        achados.append(achado("dica", "titulo-do-documento", "2.4.2",
                              "propriedades", "documento sem título nos metadados"))
    for i, shape in enumerate(doc.inline_shapes, 1):
        if not (shape._inline.docPr.get("descr") or "").strip():
            achados.append(achado("erro", "alt-text-ausente", "1.1.1",
                                  f"imagem inline nº {i}",
                                  "sem texto alternativo (descr vazio); se decorativa, marcar como decorativa"))
    for t, tabela in enumerate(doc.tables, 1):
        tem_header = any(
            linha._tr.xpath('.//*[local-name()="tblHeader"]')
            for linha in tabela.rows[:1])
        if not tem_header:
            achados.append(achado("aviso", "tabela-sem-cabecalho", "1.3.1",
                                  f"tabela nº {t}",
                                  "primeira linha não marcada como linha de cabeçalho (Repetir como cabeçalho)"))
    niveis = [int(p.style.name.split()[-1]) for p in doc.paragraphs
              if p.style.name.startswith("Heading ") and p.style.name.split()[-1].isdigit()]
    for a, b in zip(niveis, niveis[1:]):
        if b > a + 1:
            achados.append(achado("aviso", "salto-de-heading", "1.3.1",
                                  f"H{a} seguido de H{b}",
                                  "nível de título pulado quebra a navegação por leitor de tela"))
            break
    achados.append(achado("dica", "triagem-llm", "1.1.1/1.3.2",
                          "documento inteiro",
                          "qualidade de alt-text e ordem de leitura exigem julgamento: triagem LLM + confirmação humana"))
    return achados


def auditar_pptx(caminho):
    from pptx import Presentation
    prs = Presentation(caminho)
    achados = []
    for n, slide in enumerate(prs.slides, 1):
        titulo = slide.shapes.title
        if titulo is None or not (titulo.text or "").strip():
            achados.append(achado("erro", "slide-sem-titulo", "2.4.2",
                                  f"slide {n}",
                                  "todo slide precisa de título único (é o heading de navegação do leitor de tela)"))
        for shape in slide.shapes:
            if shape.shape_type == 13 and not _descr(shape._element):  # 13 = PICTURE
                achados.append(achado("erro", "alt-text-ausente", "1.1.1",
                                      f"slide {n}, imagem '{shape.name}'",
                                      "imagem sem texto alternativo (descr vazio)"))
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table._tbl
                if tbl.get("firstRow") not in ("1", "true"):
                    achados.append(achado("aviso", "tabela-sem-cabecalho", "1.3.1",
                                          f"slide {n}, tabela '{shape.name}'",
                                          "tabela sem firstRow=1 (linha de cabeçalho desligada)"))
        achados.append(achado("dica", "ordem-de-leitura", "1.3.2",
                              f"slide {n}",
                              "ordem no spTree: " + " → ".join(s.name for s in slide.shapes)
                              + " (validar se corresponde à leitura visual)"))
    return achados


def auditar_xlsx(caminho):
    import openpyxl
    wb = openpyxl.load_workbook(caminho)
    achados = []
    default_names = {"Sheet", "Sheet1", "Sheet2", "Sheet3", "Planilha1", "Planilha2", "Plan1"}
    for ws in wb.worksheets:
        if ws.title in default_names:
            achados.append(achado("aviso", "aba-sem-nome-descritivo", "2.4.2",
                                  f"aba '{ws.title}'",
                                  "renomear com nome descritivo do conteúdo"))
        if ws.merged_cells.ranges:
            achados.append(achado("aviso", "celulas-mescladas", "1.3.1",
                                  f"aba '{ws.title}'",
                                  f"{len(ws.merged_cells.ranges)} intervalo(s) mesclado(s) quebram a contagem do leitor de tela"))
        if ws["A1"].value in (None, ""):
            achados.append(achado("dica", "a1-vazia", "1.3.1",
                                  f"aba '{ws.title}'",
                                  "A1 vazia: leitores de tela começam por A1; colocar título ali"))
        if not any(getattr(ws, "tables", {}) or {}):
            achados.append(achado("dica", "dados-fora-de-tabela", "1.3.1",
                                  f"aba '{ws.title}'",
                                  "dados sem objeto Tabela (com header row) perdem semântica de cabeçalho"))
    return achados


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        sys.exit(__doc__)
    caminho = Path(args[0])
    rotas = {".docx": ("docx", auditar_docx), ".pptx": ("pptx", auditar_pptx),
             ".xlsx": ("openpyxl", auditar_xlsx)}
    if caminho.suffix.lower() not in rotas:
        sys.exit(f"formato não suportado: {caminho.suffix} (aceito: docx, pptx, xlsx)")
    mod, fn = rotas[caminho.suffix.lower()]
    try:
        achados = fn(caminho)
    except ImportError:
        sys.exit(f"extra ausente para {caminho.suffix}: rode preflight.py e instale "
                 "python-docx/python-pptx/openpyxl. Sem o extra, use o checker do "
                 "próprio Office (Revisão > Verificar Acessibilidade) e declare no relatório.")
    if "--json" in sys.argv:
        print(json.dumps({"arquivo": str(caminho), "achados": achados},
                         ensure_ascii=False, indent=2))
        return
    print(f"{caminho.name}: {len(achados)} achado(s)")
    for a in sorted(achados, key=lambda x: {"erro": 0, "aviso": 1, "dica": 2}[x["severidade"]]):
        print(f"  [{a['severidade'].upper():5}] {a['regra']} (WCAG {a['wcag']}) "
              f"· {a['local']}: {a['detalhe']}")


if __name__ == "__main__":
    main()
